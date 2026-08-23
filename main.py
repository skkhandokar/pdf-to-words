import os
import shutil
import tempfile
import subprocess
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pdf2docx import Converter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path

# FastAPI App ইনিশিয়ালাইজেশন
app = FastAPI(
    title="Shortfy PDF Converter API",
    description="High performance FastAPI backend for PDF & Office document conversions."
)

# Next.js Frontend কানেক্ট করার জন্য CORS চালুকরণ
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # Production-এ নির্দিষ্ট Next.js URL দিতে পারেন
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def home():
    return {"status": "Online", "message": "Shortfy Converter API is running smoothly!"}

@app.post("/api/convert/")
async def convert_file(
    file: UploadFile = File(...),
    conversion_type: str = Form(...) # 'pdf-to-word', 'pdf-to-excel', 'pdf-to-ppt', 'office-to-pdf'
):
    # ১. অস্থায়ী (Temp) ডিরেক্টরি ও ফাইল তৈরি
    temp_dir = tempfile.mkdtemp()
    input_path = os.path.join(temp_dir, file.filename)

    try:
        # আপলোড করা ফাইল ডিস্কে সেভ করা
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        output_path = ""
        download_filename = ""

        # ==========================================
        # ২. PDF to Word (.docx)
        # ==========================================
        if conversion_type == "pdf-to-word":
            output_path = os.path.join(temp_dir, "converted.docx")
            download_filename = "converted.docx"
            
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()

        # ==========================================
        # ৩. PDF to Excel (.xlsx)
        # ==========================================
        elif conversion_type == "pdf-to-excel":
            output_path = os.path.join(temp_dir, "converted.xlsx")
            download_filename = "converted.xlsx"
            
            all_tables = []
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    # ১. প্রথমে সাধারণ বর্ডারযুক্ত টেবিল খোঁজা
                    extracted = page.extract_tables()
                    
                    if extracted:
                        for table in extracted:
                            if table and len(table) > 1:
                                df = pd.DataFrame(table[1:], columns=table[0])
                                all_tables.append(df)
                    else:
                        # ২. যদি বর্ডার না থাকে, তবে টেক্সটের পজিশন বা লাইন ধরে এক্সট্রাক্ট করা (Implicit Table)
                        text_lines = page.extract_text()
                        if text_lines:
                            lines_data = []
                            for line in text_lines.split('\n'):
                                # স্পেস বা ট্যাব অনুযায়ী কলাম আলাদা করা
                                cols = line.split()
                                if cols:
                                    lines_data.append(cols)
                            
                            if len(lines_data) > 1:
                                # প্রথম লাইনকে হেডার এবং বাকিগুলোকে রো হিসেবে ধরা
                                max_cols = max(len(row) for row in lines_data)
                                normalized_rows = [row + [''] * (max_cols - len(row)) for row in lines_data]
                                df = pd.DataFrame(normalized_rows[1:], columns=normalized_rows[0])
                                all_tables.append(df)
            
            if all_tables:
                final_df = pd.concat(all_tables, ignore_index=True)
                final_df.to_excel(output_path, index=False)
            else:
                pd.DataFrame([{"Message": "No tabular data detected in this PDF"}]).to_excel(output_path, index=False)

        # ==========================================
        # ৪. PDF to PowerPoint (.pptx)
        # ==========================================
        elif conversion_type == "pdf-to-ppt":
            output_path = os.path.join(temp_dir, "converted.pptx")
            download_filename = "converted.pptx"
            
            images = convert_from_path(input_path)
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]

            for i, img in enumerate(images):
                img_path = os.path.join(temp_dir, f"page_{i}.png")
                img.save(img_path, "PNG")
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

            prs.save(output_path)

        # ==========================================
        # ৫. Office Documents to PDF (LibreOffice CLI)
        # ==========================================
        elif conversion_type == "office-to-pdf":
            filename_without_ext = os.path.splitext(file.filename)[0]
            output_path = os.path.join(temp_dir, f"{filename_without_ext}.pdf")
            download_filename = f"{filename_without_ext}.pdf"

            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", input_path, "--outdir", temp_dir],
                check=True
            )

        else:
            raise HTTPException(status_code=400, detail="Invalid conversion_type specified.")

        # কনভার্ট হওয়া ফাইল ইউজারকে রেসপন্স হিসেবে ডাউনলোড করিয়ে দেওয়া
        return FileResponse(
            path=output_path,
            filename=download_filename,
            media_type="application/octet-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Conversion Error: {str(e)}")